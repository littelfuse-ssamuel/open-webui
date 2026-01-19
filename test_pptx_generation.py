#!/usr/bin/env python3
"""
Test script to generate a sample PPTX file for testing the PPTX artifact viewer.
This demonstrates what the code interpreter would do.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_sample_pptx():
    """Create a sample PPTX file with multiple slides"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Sample Presentation"
    subtitle.text = "Generated for Testing PPTX Artifacts\nCreated by Code Interpreter"
    
    # Slide 2: Bullet points
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Key Features"
    
    tf = body_shape.text_frame
    tf.text = "Automatic PPTX Detection"
    
    p = tf.add_paragraph()
    p.text = "Upload to Azure Blob Storage"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Same pattern as Excel artifacts"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Middleware integration"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "File artifact emission"
    p.level = 1
    
    # Slide 3: Table
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Sales Data"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Add table
    rows = 5
    cols = 4
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(3)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(2.0)
    
    # Header row
    table.cell(0, 0).text = "Product"
    table.cell(0, 1).text = "Q1"
    table.cell(0, 2).text = "Q2"
    table.cell(0, 3).text = "Total"
    
    # Data rows
    table.cell(1, 0).text = "Laptop"
    table.cell(1, 1).text = "$15,000"
    table.cell(1, 2).text = "$18,000"
    table.cell(1, 3).text = "$33,000"
    
    table.cell(2, 0).text = "Mouse"
    table.cell(2, 1).text = "$2,500"
    table.cell(2, 2).text = "$3,000"
    table.cell(2, 3).text = "$5,500"
    
    table.cell(3, 0).text = "Keyboard"
    table.cell(3, 1).text = "$3,500"
    table.cell(3, 2).text = "$4,000"
    table.cell(3, 3).text = "$7,500"
    
    table.cell(4, 0).text = "TOTAL"
    table.cell(4, 1).text = "$21,000"
    table.cell(4, 2).text = "$25,000"
    table.cell(4, 3).text = "$46,000"
    
    # Slide 4: Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Summary"
    
    tf = body_shape.text_frame
    tf.text = "PPTX files are now automatically detected"
    
    p = tf.add_paragraph()
    p.text = "Files are uploaded to storage provider"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Artifacts are emitted to frontend"
    p.level = 0
    
    # Save the file
    output_path = "/tmp/sample_presentation.pptx"
    prs.save(output_path)
    print(f"✅ Sample PPTX file created: {output_path}")
    print(f"📊 File size: {os.path.getsize(output_path)} bytes")
    print(f"📄 Slide count: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    create_sample_pptx()
