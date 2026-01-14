"""
PowerPoint generation utilities using python-pptx.
Generates .pptx files from structured JSON slide data.

Optimized for Littelfuse corporate template with proper layout selection,
markdown formatting support, and brand color compliance.
"""

import io
import re
import base64
import logging
from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from open_webui.env import SRC_LOG_LEVELS

log = logging.getLogger(__name__)
log.setLevel(SRC_LOG_LEVELS["MAIN"])


# ================================
# Littelfuse Brand Constants
# ================================

# Primary brand colors
LITTELFUSE_GREEN = '007E3A'      # RGB(0, 126, 58) - Primary brand green
LITTELFUSE_GREY = 'CCCCCC'       # RGB(204, 204, 204) - Secondary grey

# Secondary palette
GREEN_APPLE = '78BE20'           # RGB(120, 190, 32)
LIMA_GREEN = '48A23F'            # RGB(72, 162, 63)
BONDI_BLUE = '0093B2'            # RGB(0, 147, 178)
PACIFIC_BLUE = '00AEC7'          # RGB(0, 174, 199)

# Accent colors (use sparingly - negative values only)
ACCENT_YELLOW = 'EBB93B'         # RGB(235, 185, 59)
ACCENT_RED = 'D21500'            # RGB(210, 21, 0)

# Alternating row color for tables
TABLE_ALT_ROW = 'E8F5E9'         # Light green tint for alternating rows

# Default font (matches Littelfuse template guidelines)
DEFAULT_FONT = 'Arial'

# Layout mapping for Littelfuse template
# Based on analysis of littelfuse_template_2026.pptx
LAYOUT_MAPPING = {
    'title_slide': 0,            # 13_Section Divider - Green background title
    'section_divider': 2,        # 17_Section Divider - Blue cubes background
    'content': 1,                # Title and Content - Standard bullet/text content
    'content_nda': 10,           # NDA LFUS Title and Content - With NDA banner
    'content_internal': 13,      # Internal LFUS Title and Content
    'title_only': 4,             # NDA LFUS Title Only - For tables/custom content
    'two_column': 5,             # Two Content - Side by side comparison
    'blank': 6,                  # Blank and footer
    'ending': 16,                # Ending Slide - Closing disclaimer
}


# ================================
# Markdown Processing Functions
# ================================

def strip_markdown_to_plain(text: str) -> str:
    """
    Remove markdown formatting and return plain text.
    Handles **bold**, *italic*, __bold__, _italic_, and combinations.
    """
    if not text:
        return ''
    
    # Remove bold markers: **text** or __text__
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)
    
    # Remove italic markers: *text* or _text_
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'(?<![a-zA-Z])_(.+?)_(?![a-zA-Z])', r'\1', text)
    
    return text


def parse_markdown_segments(text: str) -> List[Tuple[str, bool, bool]]:
    """
    Parse text with markdown formatting into segments.
    
    Returns:
        List of tuples: (text, is_bold, is_italic)
    """
    if not text:
        return []
    
    segments = []
    
    # Pattern to match markdown formatting
    # Matches: **bold**, *italic*, ***bold+italic***, or plain text
    pattern = r'(\*\*\*(.+?)\*\*\*|\*\*(.+?)\*\*|\*(.+?)\*|([^*]+))'
    
    for match in re.finditer(pattern, text):
        full_match = match.group(0)
        
        if match.group(2):  # ***bold+italic***
            segments.append((match.group(2), True, True))
        elif match.group(3):  # **bold**
            segments.append((match.group(3), True, False))
        elif match.group(4):  # *italic*
            segments.append((match.group(4), False, True))
        elif match.group(5):  # plain text
            if match.group(5).strip():  # Skip empty segments
                segments.append((match.group(5), False, False))
    
    # If no segments found, return the whole text as plain
    if not segments:
        segments = [(text, False, False)]
    
    return segments


def apply_formatted_runs(paragraph, text: str, base_size: int = 18, font_name: str = DEFAULT_FONT):
    """
    Apply text with markdown formatting to a paragraph using runs.
    
    Args:
        paragraph: python-pptx paragraph object
        text: Text potentially containing markdown formatting
        base_size: Base font size in points
        font_name: Font family name
    """
    # Clear existing text
    paragraph.clear()
    
    segments = parse_markdown_segments(text)
    
    for segment_text, is_bold, is_italic in segments:
        run = paragraph.add_run()
        run.text = segment_text
        run.font.size = Pt(base_size)
        run.font.name = font_name
        run.font.bold = is_bold
        run.font.italic = is_italic


# ================================
# Color and Formatting Utilities
# ================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return RGBColor(0, 0, 0)  # Default to black
    try:
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    except ValueError:
        return RGBColor(0, 0, 0)


def format_paragraph(paragraph, size: int = 18, bold: bool = False, font_name: str = DEFAULT_FONT):
    """Apply formatting to a paragraph."""
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.name = font_name


def format_text_frame(text_frame, size: int = 18, bold: bool = False, font_name: str = DEFAULT_FONT):
    """Apply formatting to all paragraphs in a text frame."""
    for paragraph in text_frame.paragraphs:
        format_paragraph(paragraph, size, bold, font_name)


def format_cell(cell, bold: bool = False, bg_color: Optional[str] = None, font_size: int = 14):
    """Format a table cell."""
    text_frame = cell.text_frame
    if text_frame.paragraphs:
        text_frame.paragraphs[0].font.size = Pt(font_size)
        text_frame.paragraphs[0].font.bold = bold
        text_frame.paragraphs[0].font.name = DEFAULT_FONT
    
    if bg_color:
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)


# ================================
# Content Addition Functions
# ================================

def add_content_to_slide(slide, content_items: List[Dict[str, Any]], prs: Presentation):
    """Add various content types to a slide with markdown support."""
    # Starting position for content (below title area)
    top = Inches(1.8)
    left = Inches(0.75)
    width = Inches(8.5)
    slide_height = prs.slide_height
    
    for item in content_items:
        item_type = item.get('type', '')
        
        if item_type == 'text':
            text = item.get('text', '')
            textbox = slide.shapes.add_textbox(left, top, width, Inches(1))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Apply formatted text with markdown support
            p = text_frame.paragraphs[0]
            apply_formatted_runs(p, text, base_size=18)
            
            top += Inches(0.8)
        
        elif item_type == 'bullet':
            items = item.get('items', [])
            if not items:
                continue
            
            # Calculate height based on number of items
            height = Inches(len(items) * 0.45 + 0.2)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            for i, bullet_text in enumerate(items):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                # Use proper bullet formatting instead of manual bullet character
                p.level = 0
                
                # Apply formatted text with markdown support
                apply_formatted_runs(p, bullet_text, base_size=16)
                
                # Set bullet properties
                p.space_after = Pt(6)
                
                # Add bullet character via low-level XML
                # This ensures bullets display correctly
                from pptx.oxml.ns import qn
                from lxml import etree
                
                pPr = p._p.get_or_add_pPr()
                # Remove any existing bullet settings
                for child in pPr.findall(qn('a:buNone')):
                    pPr.remove(child)
                # Add bullet character element
                buChar = etree.SubElement(pPr, qn('a:buChar'))
                buChar.set('char', '•')
            
            top += height + Inches(0.2)
        
        elif item_type == 'image':
            src = item.get('src', '')
            if not src:
                continue
            
            try:
                if src.startswith('data:'):
                    # Base64 encoded image
                    # Format: data:image/png;base64,XXXX
                    header, data = src.split(',', 1)
                    image_data = base64.b64decode(data)
                    image_stream = io.BytesIO(image_data)
                    
                    # Add image with max width of 6 inches, maintaining aspect ratio
                    pic = slide.shapes.add_picture(image_stream, left, top, width=Inches(6))
                    
                    # Calculate height based on aspect ratio
                    image_height = pic.height
                    top += Inches(image_height.inches + 0.3)
                else:
                    # URL-based image - skip for now (would need async download)
                    log.warning(f"URL-based images not yet supported: {src[:50]}...")
                    continue
            except Exception as e:
                log.error(f"Error adding image: {e}")
                continue
        
        elif item_type == 'table':
            headers = item.get('headers', [])
            rows = item.get('rows', [])
            
            if not headers:
                continue
            
            num_rows = len(rows) + 1  # +1 for header
            num_cols = len(headers)
            
            # Calculate table dimensions
            table_height = Inches(num_rows * 0.4 + 0.2)
            
            table_shape = slide.shapes.add_table(
                num_rows, num_cols,
                left, top,
                width, table_height
            )
            table = table_shape.table
            
            # Style the table with Littelfuse brand colors
            # Add headers with Littelfuse Green
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                # Strip markdown from header text
                cell.text = strip_markdown_to_plain(str(header))
                format_cell(cell, bold=True, bg_color=LITTELFUSE_GREEN, font_size=12)
                # White text on green background
                if cell.text_frame.paragraphs:
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            # Add data rows
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < num_cols:  # Ensure we don't exceed column count
                        cell = table.cell(row_idx + 1, col_idx)
                        # Strip markdown from cell text
                        cell.text = strip_markdown_to_plain(str(cell_data)) if cell_data is not None else ''
                        format_cell(cell, font_size=11)
                        # Alternating row colors with light green tint
                        if row_idx % 2 == 1:
                            format_cell(cell, bg_color=TABLE_ALT_ROW, font_size=11)
            
            top += table_height + Inches(0.3)
        
        # Safety check - don't overflow the slide
        if top > Inches(6.5):
            log.warning("Content may overflow slide boundaries")
            break


def get_layout_for_content(slide_info: Dict[str, Any], slide_index: int, total_slides: int, has_template: bool) -> int:
    """
    Determine the best layout index based on content type and slide position.
    
    Args:
        slide_info: Slide data dictionary
        slide_index: Position of slide in presentation (0-indexed)
        total_slides: Total number of slides
        has_template: Whether using Littelfuse template
        
    Returns:
        Layout index to use
    """
    title = slide_info.get('title', '')
    content = slide_info.get('content', [])
    slide_type = slide_info.get('type', '')  # Optional explicit type
    
    # Handle explicit slide type if provided
    if slide_type:
        if slide_type in LAYOUT_MAPPING:
            return LAYOUT_MAPPING[slide_type]
    
    # First slide is typically a title slide
    if slide_index == 0:
        return LAYOUT_MAPPING.get('title_slide', 0)
    
    # Last slide can be ending slide (if explicitly marked or no content)
    if slide_index == total_slides - 1 and not content and 'ending' in str(slide_info.get('title', '')).lower():
        return LAYOUT_MAPPING.get('ending', 16)
    
    # Analyze content to determine best layout
    has_table = any(item.get('type') == 'table' for item in content)
    has_bullets = any(item.get('type') == 'bullet' for item in content)
    has_image = any(item.get('type') == 'image' for item in content)
    has_text = any(item.get('type') == 'text' for item in content)
    
    # Table slides work best with Title Only layout (more space)
    if has_table:
        return LAYOUT_MAPPING.get('title_only', 4)
    
    # Standard content layout for bullets and text
    if has_bullets or has_text:
        return LAYOUT_MAPPING.get('content', 1)
    
    # Title only if no content
    if title and not content:
        return LAYOUT_MAPPING.get('title_only', 4)
    
    # Blank for anything else
    return LAYOUT_MAPPING.get('blank', 6)


# ================================
# Slide Creation Functions
# ================================

def create_slide(prs: Presentation, slide_info: Dict[str, Any], slide_index: int = 0, total_slides: int = 1, has_template: bool = False) -> Any:
    """
    Create a single slide based on slide_info structure.
    
    Args:
        prs: Presentation object
        slide_info: Dictionary with slide data
        slide_index: Position of this slide (0-indexed)
        total_slides: Total number of slides being created
        has_template: Whether using Littelfuse template
        
    Returns:
        Created slide object
    """
    title = slide_info.get('title', '')
    content = slide_info.get('content', [])
    background_color = slide_info.get('backgroundColor')
    notes = slide_info.get('notes', '')
    
    # Get optimal layout for this content
    layout_index = get_layout_for_content(slide_info, slide_index, total_slides, has_template)
    
    # Ensure layout index is valid
    if layout_index >= len(prs.slide_layouts):
        log.warning(f"Layout index {layout_index} out of range, falling back to 0")
        layout_index = 0
    
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    log.debug(f"Created slide {slide_index} with layout {layout_index} ({slide_layout.name})")
    
    # Set background color if specified (usually not needed with template)
    if background_color:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(background_color)
    
    # Add title with markdown support
    if title:
        # Strip markdown for title (titles shouldn't have inline formatting)
        clean_title = strip_markdown_to_plain(title)
        
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = clean_title
            format_text_frame(title_shape.text_frame, size=32, bold=True)
        else:
            # Create title textbox manually if no placeholder
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(9), Inches(1)
            )
            title_frame = title_box.text_frame
            title_frame.text = clean_title
            format_text_frame(title_frame, size=32, bold=True)
    
    # Add content
    if content:
        add_content_to_slide(slide, content, prs)
    
    # Add speaker notes with markdown stripped
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = strip_markdown_to_plain(notes)
    
    return slide


def clear_template_slides(prs: Presentation) -> None:
    """
    Remove all existing slides from a template presentation.
    This allows using the template's layouts/masters without sample content.
    """
    # Delete slides in reverse order to avoid index shifting issues
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    log.info("Cleared all template slides")


# ================================
# Main Generation Function
# ================================

def generate_pptx_from_data(
    slide_data: Dict[str, Any],
    template_path: Optional[str] = None,
    clear_template: bool = True
) -> bytes:
    """
    Generate a PowerPoint file from structured slide data.
    
    Args:
        slide_data: Dictionary containing presentation structure:
            {
                "title": "Presentation Title",
                "slides": [
                    {
                        "title": "Slide Title",
                        "type": "content",  # Optional: title_slide, content, table, ending
                        "backgroundColor": "#FFFFFF",
                        "content": [
                            {"type": "text", "text": "Some **bold** text"},
                            {"type": "bullet", "items": ["Item 1", "**Bold** item"]},
                            {"type": "table", "headers": [...], "rows": [...]},
                            {"type": "image", "src": "data:image/png;base64,..."}
                        ],
                        "notes": "Speaker notes"
                    }
                ]
            }
        template_path: Optional path to a .pptx template file
        clear_template: If True, remove sample slides from template (default True)
        
    Returns:
        bytes: PPTX file content as bytes
    """
    has_template = False
    
    # Create presentation from template or blank
    if template_path and Path(template_path).exists():
        try:
            prs = Presentation(template_path)
            has_template = True
            log.info(f"Using template: {template_path}")
            
            # Clear existing sample slides from template
            if clear_template:
                clear_template_slides(prs)
                
        except Exception as e:
            log.warning(f"Failed to load template: {e}. Using blank presentation.")
            prs = Presentation()
    else:
        prs = Presentation()
    
    # Don't override template dimensions - use what's in the template
    if not has_template:
        # Set slide dimensions (16:9 aspect ratio) only for blank presentations
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # Get slides from data
    slides = slide_data.get('slides', [])
    
    if not slides:
        # Create at least one title slide if no slides provided
        title = slide_data.get('title', 'Untitled Presentation')
        slides = [{'title': title, 'content': [], 'type': 'title_slide'}]
    
    total_slides = len(slides)
    
    # Create each slide
    for idx, slide_info in enumerate(slides):
        create_slide(prs, slide_info, slide_index=idx, total_slides=total_slides, has_template=has_template)
    
    # Save to bytes buffer
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    log.info(f"Generated presentation with {total_slides} slides")
    
    return pptx_bytes.getvalue()


# ================================
# Utility Functions
# ================================

def get_presentation_info(pptx_path: str) -> Dict[str, Any]:
    """
    Extract metadata from an existing presentation.
    
    Args:
        pptx_path: Path to the .pptx file
        
    Returns:
        Dictionary with presentation metadata
    """
    try:
        prs = Presentation(pptx_path)
        
        layouts = []
        for i, layout in enumerate(prs.slide_layouts):
            layouts.append({
                'index': i,
                'name': layout.name
            })
        
        return {
            'slide_count': len(prs.slides),
            'width_inches': prs.slide_width.inches,
            'height_inches': prs.slide_height.inches,
            'layouts': layouts
        }
    except Exception as e:
        log.error(f"Error reading presentation info: {e}")
        return {'error': str(e)}