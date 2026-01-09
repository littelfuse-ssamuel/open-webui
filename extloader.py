import fitz  # PyMuPDF
import base64
import logging
import os
import tempfile
import time
import hashlib
import json
import re
import unicodedata
from typing import Dict, List, Tuple, Any, Optional
from io import BytesIO
from pathlib import Path

from fastapi import FastAPI, Request, HTTPException, status
from fastapi.responses import JSONResponse
from dataclasses import dataclass, asdict
from datetime import datetime

# =============================================================================
# AZURE IMPORTS (The new brain)
# =============================================================================
try:
    from azure.core.credentials import AzureKeyCredential
    from azure.ai.documentintelligence import DocumentIntelligenceClient
    from azure.ai.documentintelligence.models import AnalyzeResult, AnalyzeDocumentRequest
    AZURE_AVAILABLE = True
except ImportError:
    AZURE_AVAILABLE = False

# =============================================================================
# LEGACY / FAST PROCESSOR IMPORTS
# =============================================================================
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False

# Image processing
try:
    from PIL import Image
    import io
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

# Config
DEFAULT_OUTPUT_FORMAT = "json"  # or "markdown"

# Image processing settings
AUTO_COMPRESS_IMAGES = os.getenv("AUTO_COMPRESS_IMAGES", "true").lower() == "true"
DEFAULT_COMPRESSION_WIDTH = int(os.getenv("FILE_IMAGE_COMPRESSION_WIDTH", "1024"))
DEFAULT_COMPRESSION_HEIGHT = int(os.getenv("FILE_IMAGE_COMPRESSION_HEIGHT", "1024"))
COMPRESSION_QUALITY = int(os.getenv("IMAGE_COMPRESSION_QUALITY", "85"))
MIN_IMAGE_WIDTH = int(os.getenv("MIN_IMAGE_WIDTH", "32"))
MIN_IMAGE_HEIGHT = int(os.getenv("MIN_IMAGE_HEIGHT", "32"))
MIN_IMAGE_SIZE_BYTES = int(os.getenv("MIN_IMAGE_SIZE_BYTES", "1024"))

# Debug output settings
EXTLOADER_DEBUG_OUTPUT = os.getenv("EXTLOADER_DEBUG_OUTPUT", "false").lower() == "true"
EXTLOADER_DEBUG_PATH = os.getenv("EXTLOADER_DEBUG_PATH", "/app/backend/data/extloader_debug")

# Azure Configuration
AZURE_ENDPOINT = os.getenv("DOCUMENT_INTELLIGENCE_ENDPOINT", "")
AZURE_KEY = os.getenv("DOCUMENT_INTELLIGENCE_KEY", "")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize the FastAPI application
app = FastAPI(
    title="Enhanced Content Processing Engine (Azure + PyMuPDF)",
    description="An API to extract text, tables, layout, and metadata using Azure Document Intelligence & PyMuPDF.",
    version="3.1.0"
)

# =============================================================================
# AZURE CLIENT HELPERS
# =============================================================================

def get_azure_client():
    if not AZURE_AVAILABLE:
        return None
    if not AZURE_ENDPOINT or not AZURE_KEY:
        logger.error("Azure credentials missing. Set AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT and AZURE_DOCUMENT_INTELLIGENCE_KEY.")
        return None
    
    return DocumentIntelligenceClient(
        endpoint=AZURE_ENDPOINT, 
        credential=AzureKeyCredential(AZURE_KEY)
    )

# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class DocumentMetadata:
    """Document-level metadata."""
    title: Optional[str] = None
    author: Optional[str] = None
    subject: Optional[str] = None
    keywords: Optional[str] = None
    creator: Optional[str] = None
    producer: Optional[str] = None
    creation_date: Optional[str] = None
    modification_date: Optional[str] = None
    format: Optional[str] = None
    encryption: Optional[str] = None
    page_count: int = 0
    
    def to_dict(self) -> dict:
        return {k: v for k, v in asdict(self).items() if v is not None}

@dataclass
class TableCell:
    """A single cell in an extracted table."""
    row_index: int
    col_index: int
    content: str
    row_span: int = 1
    col_span: int = 1
    is_header: bool = False

@dataclass
class ExtractedTable:
    """A table extracted from a document page."""
    page_number: int
    table_index: int
    row_count: int
    col_count: int
    cells: List[TableCell]
    bbox: Tuple[float, float, float, float]  # x0, y0, x1, y1
    markdown: str = ""
    header_external: bool = False
    
    def to_dict(self) -> dict:
        return {
            "page_number": self.page_number,
            "table_index": self.table_index,
            "row_count": self.row_count,
            "col_count": self.col_count,
            "cells": [asdict(c) for c in self.cells],
            "bbox": self.bbox,
            "markdown": self.markdown,
            "header_external": self.header_external
        }

# =============================================================================
# TEXT UTILS & RECONSTRUCTION
# =============================================================================

def normalize_text_encoding(text: str) -> str:
    """Normalize text encoding to handle unicode issues."""
    if not text:
        return ""
    try:
        normalized = unicodedata.normalize('NFKC', text)
        replacements = {
            '\u00a0': ' ', '\xa0': ' ', '\u2013': '-', '\u2014': '--',
            '\u2018': "'", '\u2019': "'", '\u201c': '"', '\u201d': '"',
            '\u2022': '*', '\u2026': '...',
        }
        for unicode_char, replacement in replacements.items():
            normalized = normalized.replace(unicode_char, replacement)
        return normalized.encode('utf-8', errors='ignore').decode('utf-8', errors='ignore')
    except Exception:
        return text

def azure_table_to_markdown(table) -> str:
    """
    Converts an Azure Table object into a Markdown string.
    This is required to embed tables into the page text.
    """
    if not table.row_count or not table.column_count:
        return ""

    # Initialize grid
    grid = [["" for _ in range(table.column_count)] for _ in range(table.row_count)]
    
    # Fill grid
    for cell in table.cells:
        r = cell.row_index
        c = cell.column_index
        content = normalize_text_encoding(cell.content).replace("\n", " ")
        # Handle spans roughly by filling the top-left (rendering limits apply)
        if r < table.row_count and c < table.column_count:
            grid[r][c] = content

    # Build Markdown
    md_lines = []
    
    # Header row
    header_row = grid[0]
    md_lines.append("| " + " | ".join(header_row) + " |")
    md_lines.append("| " + " | ".join(["---"] * len(header_row)) + " |")
    
    # Data rows
    for row in grid[1:]:
        md_lines.append("| " + " | ".join(row) + " |")
        
    return "\n".join(md_lines)

def reconstruct_pages_from_azure(result: Any) -> Tuple[List[str], str]:
    """
    Reconstructs page-by-page text from Azure result structure.
    This ensures strict page splitting by mapping paragraphs and tables 
    to their page numbers and sorting by reading order (offset).
    """
    if not result:
        return [], ""
        
    # 1. Bucketize content by page
    # Maps page_number (1-based) -> list of (offset, content_type, content_obj)
    page_map: Dict[int, List[Tuple[int, str, Any]]] = {}
    total_pages = len(result.pages) if result.pages else 0
    
    # Initialize buckets
    for i in range(1, total_pages + 1):
        page_map[i] = []

    # 2. Process Paragraphs (Text)
    if result.paragraphs:
        for p in result.paragraphs:
            if not p.bounding_regions: continue
            
            page_num = p.bounding_regions[0].page_number
            offset = p.spans[0].offset if p.spans else 0
            
            # Determine role for markdown formatting
            text = p.content
            role = getattr(p, "role", None)
            
            if role == "title":
                text = f"# {text}"
            elif role == "sectionHeading":
                text = f"## {text}"
                
            if page_num in page_map:
                page_map[page_num].append((offset, "text", text))

    # 3. Process Tables
    # We treat tables as a single block inserted at their starting offset
    if result.tables:
        for t in result.tables:
            if not t.bounding_regions: continue
            
            page_num = t.bounding_regions[0].page_number
            offset = t.spans[0].offset if t.spans else 0
            
            # Convert table to markdown immediately for storage
            table_md = azure_table_to_markdown(t)
            
            if page_num in page_map:
                page_map[page_num].append((offset, "table", table_md))

    # 4. Construct Final Pages
    final_pages = []
    full_content = []
    
    for i in range(1, total_pages + 1):
        items = page_map[i]
        # Sort by offset to preserve reading order
        items.sort(key=lambda x: x[0])
        
        # Join content with newlines
        page_text = "\n\n".join([item[2] for item in items])
        final_pages.append(page_text)
        full_content.append(page_text)
        
    return final_pages, "\n\n".join(full_content)

def write_debug_output(filename: str, page_content: str, full_result: dict) -> None:
    if not EXTLOADER_DEBUG_OUTPUT:
        return
    try:
        os.makedirs(EXTLOADER_DEBUG_PATH, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_filename = re.sub(r'[^\w\-.]', '_', filename)[:50]
        prefix = f"{timestamp}_{safe_filename}"
        
        with open(os.path.join(EXTLOADER_DEBUG_PATH, f"{prefix}_content.md"), 'w', encoding='utf-8') as f:
            f.write(page_content)
        
        with open(os.path.join(EXTLOADER_DEBUG_PATH, f"{prefix}_full.json"), 'w', encoding='utf-8') as f:
            json.dump(full_result, f, indent=2, ensure_ascii=False, default=str)
            
        logger.info(f"Debug output written to {EXTLOADER_DEBUG_PATH}")
    except Exception as e:
        logger.warning(f"Failed to write debug output: {e}")

# =============================================================================
# METADATA & UTILS (PyMuPDF)
# =============================================================================

def extract_pdf_metadata(doc: fitz.Document) -> DocumentMetadata:
    """Extract standard metadata using PyMuPDF (Fast)."""
    metadata = DocumentMetadata()
    try:
        pdf_metadata = doc.metadata
        if pdf_metadata:
            metadata.title = pdf_metadata.get('title')
            metadata.author = pdf_metadata.get('author')
            metadata.subject = pdf_metadata.get('subject')
            metadata.keywords = pdf_metadata.get('keywords')
            metadata.creator = pdf_metadata.get('creator')
            metadata.producer = pdf_metadata.get('producer')
            metadata.format = pdf_metadata.get('format')
            metadata.encryption = pdf_metadata.get('encryption')
            
            creation_date = pdf_metadata.get('creationDate')
            if creation_date:
                metadata.creation_date = parse_pdf_date(creation_date)
            
            mod_date = pdf_metadata.get('modDate')
            if mod_date:
                metadata.modification_date = parse_pdf_date(mod_date)
        
        metadata.page_count = doc.page_count
    except Exception as e:
        logger.warning(f"Error extracting metadata: {e}")
        metadata.page_count = doc.page_count if doc else 0
    return metadata

def parse_pdf_date(date_string: str) -> Optional[str]:
    if not date_string: return None
    try:
        if date_string.startswith("D:"): date_string = date_string[2:]
        year = date_string[0:4] if len(date_string) >= 4 else None
        month = date_string[4:6] if len(date_string) >= 6 else "01"
        day = date_string[6:8] if len(date_string) >= 8 else "01"
        hour = date_string[8:10] if len(date_string) >= 10 else "00"
        minute = date_string[10:12] if len(date_string) >= 12 else "00"
        second = date_string[12:14] if len(date_string) >= 14 else "00"
        if year: return f"{year}-{month}-{day}T{hour}:{minute}:{second}"
    except Exception: pass
    return date_string

def get_table_of_contents(doc: fitz.Document) -> List[Dict]:
    toc = []
    try:
        raw_toc = doc.get_toc()
        for entry in raw_toc:
            if len(entry) >= 3:
                toc.append({
                    "level": entry[0],
                    "title": normalize_text_encoding(entry[1]),
                    "page": entry[2]
                })
    except Exception: pass
    return toc

# =============================================================================
# HYBRID PDF PROCESSING
# =============================================================================

def process_pdf_enhanced(
    file_bytes: bytes, 
    filename: str,
    extract_images_flag: str = "false",
    extract_tables: bool = True,
    detect_layout: bool = True,
    extract_metadata: bool = True,
    exclude_headers_footers: bool = False,
    output_format: str = "json"
) -> Dict[str, Any]:
    """
    Enhanced PDF processing using Azure for content/layout and PyMuPDF for assets/metadata.
    """
    start_time = time.time()
    
    result = {
        "page_content": "",
        "pages": [],
        "tables": [],
        "metadata": {
            "source": filename,
            "processing_status": "pending",
            "engine": "azure_hybrid"
        },
        "images": [],
        "document_metadata": {},
        "table_of_contents": [],
        "layout_info": []
    }

    # -------------------------------------------------------------------------
    # PHASE 1: PyMuPDF (Fast Pass)
    # Extract Metadata, TOC, and Raw Images
    # -------------------------------------------------------------------------
    try:
        with fitz.open(stream=file_bytes, filetype="pdf") as doc_fitz:
            if extract_metadata:
                result["document_metadata"] = extract_pdf_metadata(doc_fitz).to_dict()
                result["table_of_contents"] = get_table_of_contents(doc_fitz)
                
                # Basic layout info 
                for page_num in range(len(doc_fitz)):
                    page = doc_fitz[page_num]
                    result["layout_info"].append({
                        "page_number": page_num + 1,
                        "width": page.rect.width,
                        "height": page.rect.height,
                        "column_count": 1,
                        "has_header": False, 
                        "has_footer": False
                    })

            if extract_images_flag.lower() == "true":
                result["images"] = process_images_efficiently(doc_fitz, extract_images_flag, filename)
    except Exception as e:
        logger.error(f"PyMuPDF Phase failed: {e}")
        # Continue to Azure phase even if fitz fails (unlikely)

    # -------------------------------------------------------------------------
    # PHASE 2: Azure Document Intelligence (Deep Pass)
    # Extract Layout-Aware Text and Semantic Tables
    # -------------------------------------------------------------------------
    client = get_azure_client()
    if not client:
        raise HTTPException(status_code=500, detail="Azure Document Intelligence client not available.")

    try:
        logger.info(f"Starting Azure analysis for {filename}...")
        
        # Analyze Document (SDK v1.0.0+ requires body parameter)
        poller = client.begin_analyze_document(
            model_id="prebuilt-layout",
            body=AnalyzeDocumentRequest(bytes_source=file_bytes)
        )
        analyze_result: AnalyzeResult = poller.result()
        
        # 1. Reconstruct Pages & Text
        # We manually stitch paragraphs and tables together to respect page boundaries
        pages_list, full_text = reconstruct_pages_from_azure(analyze_result)
        result["pages"] = pages_list
        result["page_content"] = full_text
        
        # 2. Extract Tables (Structured)
        if extract_tables and analyze_result.tables:
            for i, table in enumerate(analyze_result.tables):
                try:
                    # Convert Azure Cells to Internal Schema
                    cells = []
                    for cell in table.cells:
                        is_header = getattr(cell, "kind", "") == "columnHeader"
                        cells.append(TableCell(
                            row_index=cell.row_index,
                            col_index=cell.column_index,
                            content=normalize_text_encoding(cell.content),
                            row_span=cell.row_span if cell.row_span else 1,
                            col_span=cell.column_span if cell.column_span else 1,
                            is_header=is_header
                        ))
                    
                    # Page number (default to 1, or the first bounding region)
                    page_no = 1
                    if table.bounding_regions:
                        page_no = table.bounding_regions[0].page_number
                        
                    # Calculate simple bbox from polygon
                    # Azure polygon is [x1, y1, x2, y2, x3, y3, x4, y4]
                    bbox = (0,0,0,0)
                    if table.bounding_regions and table.bounding_regions[0].polygon:
                        poly = table.bounding_regions[0].polygon
                        xs = poly[0::2]
                        ys = poly[1::2]
                        bbox = (min(xs), min(ys), max(xs), max(ys))

                    # Generate Markdown for this specific table
                    table_md = azure_table_to_markdown(table)

                    result["tables"].append(ExtractedTable(
                        page_number=page_no,
                        table_index=i,
                        row_count=table.row_count,
                        col_count=table.column_count,
                        cells=cells,
                        bbox=bbox,
                        markdown=table_md,
                        header_external=False
                    ).to_dict())
                except Exception as e:
                    logger.warning(f"Error converting Azure table {i}: {e}")
                    continue

    except Exception as e:
        logger.error(f"Azure analysis failed: {e}")
        raise HTTPException(status_code=500, detail=f"Azure processing failed: {str(e)}")

    # Finalize Metadata
    result["metadata"]["total_pages"] = result["document_metadata"].get("page_count", 0)
    result["metadata"]["tables_extracted"] = len(result["tables"])
    result["metadata"]["images_extracted"] = len(result["images"])
    result["metadata"]["processing_time_ms"] = int((time.time() - start_time) * 1000)
    result["metadata"]["processing_status"] = "completed"

    # Write debug output
    write_debug_output(filename, result["page_content"], result)
    
    logger.info(f"Hybrid processing complete: {len(result['tables'])} tables, {len(result['images'])} images")
    return result

# =============================================================================
# IMAGE PROCESSING (PyMuPDF - Efficient)
# =============================================================================

def compress_image_for_efficiency(image_bytes: bytes, image_ext: str) -> tuple[bytes, bool, tuple[int, int]]:
    """Compress image for processing efficiency if auto-compression is enabled."""
    if not PILLOW_AVAILABLE:
        return image_bytes, False, (0, 0)
    
    try:
        image = Image.open(io.BytesIO(image_bytes))
        original_dimensions = (image.width, image.height)
        
        if not AUTO_COMPRESS_IMAGES:
            return image_bytes, False, original_dimensions
        
        if image.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', image.size, (255, 255, 255))
            if image.mode == 'P':
                image = image.convert('RGBA')
            background.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
            image = background
        
        if image.width > DEFAULT_COMPRESSION_WIDTH or image.height > DEFAULT_COMPRESSION_HEIGHT:
            image.thumbnail((DEFAULT_COMPRESSION_WIDTH, DEFAULT_COMPRESSION_HEIGHT), Image.Resampling.LANCZOS)
        
        output = io.BytesIO()
        image.save(output, format='JPEG', quality=COMPRESSION_QUALITY, optimize=True)
        compressed_bytes = output.getvalue()
        
        if len(compressed_bytes) < len(image_bytes):
            return compressed_bytes, True, original_dimensions
        else:
            return image_bytes, False, original_dimensions
            
    except Exception as e:
        logger.warning(f"Image compression failed, using original: {e}")
        return image_bytes, False, (0, 0)

def process_images_efficiently(doc, extract_images_flag: str, filename: str) -> list:
    """Process PDF images with efficiency optimizations, filtering, and deduplication."""
    base64_images = []
    
    if extract_images_flag != 'true':
        return base64_images
    
    logger.info("Starting efficient image extraction...")
    
    seen_hashes = set()
    seen_xrefs = set()
    images_processed = 0
    images_filtered = 0
    images_duplicated = 0
    
    try:
        for page_num in range(len(doc)):
            image_list = doc.get_page_images(page_num, full=True)
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    
                    if xref in seen_xrefs:
                        images_duplicated += 1
                        continue
                    
                    seen_xrefs.add(xref)
                    
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    if len(image_bytes) < MIN_IMAGE_SIZE_BYTES:
                        images_filtered += 1
                        continue
                    
                    processed_bytes, was_compressed, dimensions = compress_image_for_efficiency(image_bytes, image_ext)
                    width, height = dimensions
                    
                    if width > 0 and height > 0:
                        if width < MIN_IMAGE_WIDTH or height < MIN_IMAGE_HEIGHT:
                            images_filtered += 1
                            continue
                    
                    image_hash = hashlib.md5(processed_bytes).hexdigest()
                    if image_hash in seen_hashes:
                        images_duplicated += 1
                        continue
                    
                    seen_hashes.add(image_hash)
                    
                    encoded_string = base64.b64encode(processed_bytes).decode("utf-8")
                    
                    if was_compressed:
                        data_uri = f"data:image/jpeg;base64,{encoded_string}"
                    else:
                        data_uri = f"data:image/{image_ext};base64,{encoded_string}"
                    
                    base64_images.append(data_uri)
                    images_processed += 1
                    
                except Exception as e:
                    logger.warning(f"Failed to process image on page {page_num + 1}: {e}")
                    continue
        
        logger.info(f"Image processing: {images_processed} extracted, {images_filtered} filtered, {images_duplicated} duplicates")
        
    except Exception as e:
        logger.error(f"Error during image processing: {e}")
    
    return base64_images

# =============================================================================
# OFFICE DOCUMENT PROCESSORS (DOCX, XLSX, PPTX)
# =============================================================================

def extract_docx_structure_azure_format(doc, filename: str, output_format: str = "json") -> str:
    """Extract structured content from DOCX in Azure Document Intelligence format."""
    
    full_content = ""
    content_elements = []
    element_id = 0
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            start_offset = len(full_content)
            text_content = normalize_text_encoding(paragraph.text.strip())
            
            style_name = paragraph.style.name if paragraph.style else ""
            if "Heading" in style_name:
                level = 1
                if "Heading 2" in style_name: level = 2
                elif "Heading 3" in style_name: level = 3
                elif "Heading 4" in style_name: level = 4
                elif "Heading 5" in style_name: level = 5
                elif "Heading 6" in style_name: level = 6
                
                element_type = "title"
                role = f"heading{level}"
            else:
                element_type = "paragraph"
                role = "paragraph"
            
            content_elements.append({
                "id": f"element_{element_id}",
                "kind": element_type,
                "role": role,
                "content": text_content,
                "boundingRegions": [{"pageNumber": 1}],
                "spans": [{"offset": start_offset, "length": len(text_content)}]
            })
            
            full_content += text_content + "\n\n"
            element_id += 1
    
    # Process tables
    table_id = 0
    for table in doc.tables:
        start_offset = len(full_content)
        table_content = ""
        
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = normalize_text_encoding(cell.text.strip())
                row_data.append(cell_text)
            table_data.append(row_data)
        
        for row in table_data:
            table_content += " | ".join(row) + "\n"
        
        content_elements.append({
            "id": f"table_{table_id}",
            "kind": "table",
            "rowCount": len(table_data),
            "columnCount": len(table_data[0]) if table_data else 0,
            "cells": [
                {
                    "kind": "content",
                    "rowIndex": row_idx,
                    "columnIndex": col_idx,
                    "content": cell_content,
                }
                for row_idx, row in enumerate(table_data)
                for col_idx, cell_content in enumerate(row)
                if cell_content.strip()
            ],
            "boundingRegions": [{"pageNumber": 1}],
            "spans": [{"offset": start_offset, "length": len(table_content)}]
        })
        
        full_content += table_content + "\n\n"
        table_id += 1
    
    if output_format.lower() == "markdown":
        markdown_content = ""
        for element in content_elements:
            if element["kind"] == "title":
                level = int(element["role"].replace("heading", ""))
                markdown_content += "#" * level + " " + element["content"] + "\n\n"
            elif element["kind"] == "paragraph":
                markdown_content += element["content"] + "\n\n"
            elif element["kind"] == "table":
                table_rows = []
                max_cols = element.get("columnCount", 0)
                
                for cell in element["cells"]:
                    row_idx = cell["rowIndex"]
                    col_idx = cell["columnIndex"]
                    while row_idx >= len(table_rows):
                        table_rows.append([""] * max_cols)
                    if col_idx < max_cols:
                        table_rows[row_idx][col_idx] = cell["content"]
                
                if table_rows:
                    markdown_content += "| " + " | ".join(table_rows[0]) + " |\n"
                    markdown_content += "|" + "---|" * len(table_rows[0]) + "\n"
                    for row in table_rows[1:]:
                        markdown_content += "| " + " | ".join(row) + " |\n"
                    markdown_content += "\n"
        
        return markdown_content.strip()
    else:
        azure_response = {
            "apiVersion": "2024-11-30",
            "modelId": "prebuilt-layout",
            "content": full_content.strip(),
            "pages": [{"pageNumber": 1, "width": 8.5, "height": 11, "unit": "inch"}],
            "paragraphs": [e for e in content_elements if e["kind"] in ["paragraph", "title"]],
            "tables": [e for e in content_elements if e["kind"] == "table"],
        }
        return json.dumps(azure_response, indent=2)

def process_docx(file_bytes: bytes, filename: str) -> tuple[str, int]:
    if not DOCX_AVAILABLE:
        raise HTTPException(status_code=500, detail="python-docx not available")
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        doc = DocxDocument(temp_file_path)
        structured_content = extract_docx_structure_azure_format(doc, filename, "json")
        os.unlink(temp_file_path)
        page_count = max(1, len(structured_content) // 3000)
        return structured_content, page_count
    except Exception as e:
        if 'temp_file_path' in locals():
            try: os.unlink(temp_file_path)
            except: pass
        raise HTTPException(status_code=500, detail=f"Failed to process DOCX: {e}")

def process_xlsx(file_bytes: bytes, filename: str) -> tuple[str, int]:
    if not OPENPYXL_AVAILABLE:
        raise HTTPException(status_code=500, detail="openpyxl not available")
    try:
        file_ext = os.path.splitext(filename)[1].lower() or '.xlsx'
        if file_ext not in ['.xlsx', '.xlsm', '.xls']: file_ext = '.xlsx'
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        workbook = openpyxl.load_workbook(temp_file_path, read_only=False, data_only=False)
        full_text = ""
        sheet_count = 0
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_count += 1
            full_text += f"\n=== Sheet: {sheet_name} ===\n\n"
            row_num = 0
            for row in sheet.iter_rows():
                row_num += 1
                row_values = []
                has_content = False
                for cell in row:
                    if cell.value is not None:
                        has_content = True
                        if isinstance(cell.value, str) and cell.value.startswith('='):
                            cell_str = f"{normalize_text_encoding(cell.value)} (formula)"
                        else:
                            cell_str = normalize_text_encoding(str(cell.value))
                        row_values.append(cell_str)
                    else:
                        row_values.append("")
                if has_content:
                    while row_values and row_values[-1] == "": row_values.pop()
                    if row_values: full_text += f"[R{row_num}] " + "\t".join(row_values) + "\n"
            full_text += "\n"
        
        workbook.close()
        os.unlink(temp_file_path)
        return full_text.strip(), sheet_count
    except Exception as e:
        if 'temp_file_path' in locals():
            try: os.unlink(temp_file_path)
            except: pass
        raise HTTPException(status_code=500, detail=f"Failed to process XLSX: {e}")

def process_pptx(file_bytes: bytes, filename: str) -> tuple[str, int]:
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=500, detail="python-pptx not available")
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        prs = Presentation(temp_file_path)
        full_text = ""
        slide_count = 0
        for slide in prs.slides:
            slide_count += 1
            full_text += f"\n=== Slide {slide_count} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = normalize_text_encoding(shape.text)
                    full_text += shape_text + "\n"
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_text = normalize_text_encoding(cell.text)
                            full_text += cell_text + " "
                        full_text += "\n"
        os.unlink(temp_file_path)
        return full_text.strip(), slide_count
    except Exception as e:
        if 'temp_file_path' in locals():
            try: os.unlink(temp_file_path)
            except: pass
        raise HTTPException(status_code=500, detail=f"Failed to process PPTX: {e}")

def process_rtf(file_bytes: bytes, filename: str) -> tuple[str, int]:
    if not RTF_AVAILABLE:
        raise HTTPException(status_code=500, detail="striprtf not available")
    try:
        rtf_content = file_bytes.decode('utf-8', errors='ignore')
        plain_text = rtf_to_text(rtf_content)
        plain_text = normalize_text_encoding(plain_text)
        page_count = max(1, len(plain_text) // 3000)
        return plain_text.strip(), page_count
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to process RTF: {e}")

def process_non_pdf_fast(file_bytes: bytes, filename: str, file_ext: str) -> tuple[str, int]:
    logger.info(f"Processing {file_ext} file: {filename}")
    if file_ext == ".docx": return process_docx(file_bytes, filename)
    elif file_ext in [".xlsx", ".xlsm"]: return process_xlsx(file_bytes, filename)
    elif file_ext == ".pptx": return process_pptx(file_bytes, filename)
    elif file_ext == ".rtf": return process_rtf(file_bytes, filename)
    else: raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_ext}")

# =============================================================================
# API ENDPOINTS
# =============================================================================

@app.put("/process")
async def process_document(request: Request):
    """
    Processes an uploaded document.
    """
    start_time = time.time()
    
    file_bytes = await request.body()
    if not file_bytes:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No file content received.")

    filename = request.headers.get("x-filename", "unknown_file")
    file_ext = os.path.splitext(filename)[1].lower()
    
    # Extract headers
    extract_images_flag = request.headers.get("x-extract-images", "false").lower()
    extract_tables = request.headers.get("x-extract-tables", "true").lower() == "true"
    detect_layout = request.headers.get("x-detect-layout", "true").lower() == "true"
    extract_metadata = request.headers.get("x-extract-metadata", "true").lower() == "true"
    exclude_headers_footers = request.headers.get("x-exclude-headers-footers", "false").lower() == "true"
    output_format = request.headers.get("outputContentFormat", DEFAULT_OUTPUT_FORMAT).lower()
    
    logger.info(f"Processing: {filename}, ext={file_ext}")

    if file_ext == ".pdf":
        result = process_pdf_enhanced(
            file_bytes=file_bytes,
            filename=filename,
            extract_images_flag=extract_images_flag,
            extract_tables=extract_tables,
            detect_layout=detect_layout,
            extract_metadata=extract_metadata,
            exclude_headers_footers=exclude_headers_footers,
            output_format=output_format
        )
        return JSONResponse(content=result)
    
    else:
        try:
            if file_ext == ".docx" and DOCX_AVAILABLE:
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                    temp_file.write(file_bytes)
                    temp_file_path = temp_file.name
                doc = DocxDocument(temp_file_path)
                full_text = extract_docx_structure_azure_format(doc, filename, output_format)
                page_count = max(1, len(full_text) // 3000)
                os.unlink(temp_file_path)
            else:
                full_text, page_count = process_non_pdf_fast(file_bytes, filename, file_ext)
            
            response_payload = {
                "page_content": full_text.strip(),
                "metadata": {
                    "source": filename,
                    "page_count": page_count,
                    "images_extracted": 0,
                    "processing_status": "completed",
                    "processing_time_ms": int((time.time() - start_time) * 1000),
                    "output_format": output_format,
                },
                "images": [],
                "tables": [],
            }
            
            write_debug_output(filename, response_payload["page_content"], response_payload)
            return JSONResponse(content=response_payload)
            
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Error processing {file_ext} document: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to process document: {e}")

@app.get("/")
def read_root():
    return {
        "status": "ok",
        "message": "Enhanced Content Processing Engine (Azure Edition)",
        "version": "3.1.0",
        "engines": {
            "layout_and_text": "Azure Document Intelligence (Prebuilt Layout)",
            "images_and_metadata": "PyMuPDF"
        },
        "processors_available": {
            "azure": AZURE_AVAILABLE,
            "docx": DOCX_AVAILABLE,
            "xlsx": OPENPYXL_AVAILABLE,
            "pptx": PPTX_AVAILABLE,
            "rtf": RTF_AVAILABLE,
        }
    }

@app.get("/health")
def health_check():
    return {
        "status": "healthy",
        "azure_ready": AZURE_AVAILABLE and bool(AZURE_ENDPOINT) and bool(AZURE_KEY),
        "processors": {
            "pdf": f"Azure {'Available' if AZURE_AVAILABLE else 'Missing SDK'} + PyMuPDF",
        }
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)